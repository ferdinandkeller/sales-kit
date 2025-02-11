"""PowerPoint manipulation logic."""

import pptx

from structs import ClientOrderEstimate
from utils import get_output_path


def ppt_process_order(
    pptx_template_filename: str,
    order_estimate: ClientOrderEstimate,
) -> None:
    """Convert a client estimate into a PowerPoint presentation."""
    prs = pptx.Presentation(pptx_template_filename)

    to_replace = {
        "[CLIENT]": order_estimate.client_name,
        "[TOTAL_COST]": f"${order_estimate.total:,.2f}",
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_text = shape.text_frame.text  # type: ignore # noqa: PGH003
            for key, replacement in to_replace.items():
                if key in shape_text:
                    shape.text_frame.text = shape_text.replace(key, replacement)  # type: ignore # noqa: PGH003

    prs.save(str(get_output_path(order_estimate.client_name, "presentation.pptx")))
